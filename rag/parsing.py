"""Deterministic and LLM-assisted parsing of diverse file types."""
import os
import json
from typing import List, Tuple, Literal
from dataclasses import dataclass
from bs4 import BeautifulSoup
from transformers import AutoTokenizer
from langchain_groq import ChatGroq
from langchain_core.messages import SystemMessage
from .config import RAGConfig
from .logging_utils import get_logger, Timer

logger = get_logger("parsing")

ParsingMethod = Literal["deterministic", "llm-assisted"]


@dataclass
class ParsedDocument:
    source_path: str
    source_filename: str
    sections_text: List[Tuple[str, str]]
    parsing_method: ParsingMethod
    page_count: int


def _read_text_fallback(path: str) -> str:
    with open(path, "rb") as f:
        raw = f.read()
    try:
        return raw.decode("utf-8")
    except Exception:
        return raw.decode("latin-1", errors="ignore")


def _parse_pdf(path: str, cfg: RAGConfig) -> ParsedDocument:
    import fitz
    d = fitz.open(path)
    total_pages = d.page_count
    if total_pages > cfg.max_pages_pdf:
        logger.info("PDF exceeds page limit", extra={"event": "page_count_check", "data": {"pages": total_pages}})
        raise ValueError(f"PDF has {total_pages} pages; limit is {cfg.max_pages_pdf}")
    sections: List[Tuple[str, str]] = []
    for i in range(total_pages):
        page = d.load_page(i)
        txt = page.get_text("text") or ""
        sections.append((f"page-{i+1}", txt))
    d.close()
    return ParsedDocument(path, os.path.basename(path), sections, "deterministic", total_pages)


def _parse_docx(path: str) -> ParsedDocument:
    import docx
    d = docx.Document(path)
    txt = "\n".join(p.text for p in d.paragraphs)
    return ParsedDocument(path, os.path.basename(path), [("section-1", txt)], "deterministic", 1)


def _parse_txt(path: str) -> ParsedDocument:
    txt = _read_text_fallback(path)
    return ParsedDocument(path, os.path.basename(path), [("section-1", txt)], "deterministic", 1)


def _parse_html(path: str) -> ParsedDocument:
    html = _read_text_fallback(path)
    soup = BeautifulSoup(html, "lxml")
    txt = soup.get_text(separator="\n")
    return ParsedDocument(path, os.path.basename(path), [("section-1", txt)], "deterministic", 1)


def _parse_csv(path: str) -> ParsedDocument:
    import csv
    rows = []
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        reader = csv.reader(f)
        for r in reader:
            rows.append(", ".join(r))
    txt = "\n".join(rows)
    return ParsedDocument(path, os.path.basename(path), [("section-1", txt)], "deterministic", 1)


def _parse_json(path: str) -> ParsedDocument:
    data = json.loads(_read_text_fallback(path))
    txt = json.dumps(data, ensure_ascii=False, indent=2)
    return ParsedDocument(path, os.path.basename(path), [("section-1", txt)], "deterministic", 1)


def _parse_md_or_code(path: str) -> ParsedDocument:
    txt = _read_text_fallback(path)
    return ParsedDocument(path, os.path.basename(path), [("section-1", txt)], "deterministic", 1)


def _parse_pptx(path: str) -> ParsedDocument:
    from pptx import Presentation
    prs = Presentation(path)
    sections: List[Tuple[str, str]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shp in slide.shapes:
            if hasattr(shp, "text"):
                t = getattr(shp, "text") or ""
                if t:
                    texts.append(t)
        sections.append((f"slide-{idx}", "\n".join(texts)))
    return ParsedDocument(path, os.path.basename(path), sections, "deterministic", len(sections))


def _llm_assisted_parse(path: str, cfg: RAGConfig) -> ParsedDocument:
    api_key = os.environ.get("GROQ_API_KEY")
    if not api_key:
        raise EnvironmentError("GROQ_API_KEY is required for LLM-assisted parsing")
    tokenizer = AutoTokenizer.from_pretrained(cfg.model_name)
    content = _read_text_fallback(path)
    ids = tokenizer.encode(content)
    if len(ids) > cfg.max_context_tokens:
        ids = ids[: cfg.max_context_tokens]
    truncated = tokenizer.decode(ids)
    prompt = (
        "You are a strict parser. Extract only readable text present in the file content below.\n"
        "- Do not add or invent text.\n"
        "- If the content is not interpretable, respond with exactly: UNPARSEABLE\n"
        "Content:\n{content}"
    )
    model = ChatGroq(temperature=0.0, groq_api_key=api_key, model_name=cfg.groq_model)
    messages = [SystemMessage(content=prompt.format(content=truncated))]
    timer = Timer()
    resp = model.invoke(messages, config={"max_tokens": cfg.max_output_tokens})
    logger.info("LLM-assisted parse", extra={"event": "llm_parse", "data": {"ms": timer.elapsed_ms()}})
    parsed = resp.content.strip()
    if parsed.upper().strip() == "UNPARSEABLE" or not parsed:
        raise ValueError("File cannot be meaningfully interpreted")
    return ParsedDocument(path, os.path.basename(path), [("section-1", parsed)], "llm-assisted", 1)


def parse_file(path: str, ext: str, cfg: RAGConfig) -> ParsedDocument:
    timer = Timer()
    try:
        if ext == ".pdf":
            doc = _parse_pdf(path, cfg)
        elif ext == ".docx":
            doc = _parse_docx(path)
        elif ext in (".txt", ".log"):
            doc = _parse_txt(path)
        elif ext in (".html", ".htm"):
            doc = _parse_html(path)
        elif ext == ".csv":
            doc = _parse_csv(path)
        elif ext == ".json":
            doc = _parse_json(path)
        elif ext in (".md", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".go", ".rs"):
            doc = _parse_md_or_code(path)
        elif ext == ".pptx":
            doc = _parse_pptx(path)
        else:
            doc = _llm_assisted_parse(path, cfg)
        logger.info(
            "Parsing complete",
            extra={
                "event": "parsing",
                "data": {"method": doc.parsing_method, "sections": len(doc.sections_text), "ms": timer.elapsed_ms()},
            },
        )
        return doc
    except Exception as e:
        logger.info("Parsing failed", extra={"event": "parsing_error", "data": {"error": str(e)}})
        raise
